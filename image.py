#!/usr/bin/env python3
"""
Vision-based slide analyzer — companion to ``analyze_pptx.py``.

Pipeline per slide:
  1. Render the .pptx to PDF (LibreOffice ``soffice``) if a sibling PDF
     doesn't already exist.
  2. Rasterize each PDF page to a PNG via PyMuPDF (~200 DPI).
  3. Slice the deck with python-pptx to get the underlying chart/table
     numbers as markdown ("ground truth" for native PPTX content).
  4. For each slide, send {image + ground-truth markdown + rich JSON
     schema} to a vision-capable Foundry chat-completions deployment
     (default ``gpt-5.4``) and parse the structured response.

Outputs (next to the input deck unless --out-dir is given):
  <deck>.vision.slides.json   rich nested per-slide records
  <deck>.vision.slides.jsonl  flat one-doc-per-slide payload for Azure AI Search

Usage:
  python analyze_pptx_vision.py path/to/deck.pptx
  python analyze_pptx_vision.py path/to/deck.pptx --pdf path/to/deck.pdf
  python analyze_pptx_vision.py path/to/deck.pptx --only-slides 9,12-14
"""

from __future__ import annotations

import argparse
import base64
import json
import os
import shutil
import subprocess
import sys
import time
from pathlib import Path

import fitz  # PyMuPDF
from azure.identity import (
    AzureCliCredential,
    ChainedTokenCredential,
    DefaultAzureCredential,
    get_bearer_token_provider,
)
from dotenv import load_dotenv
from openai import AzureOpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Inlined python-pptx -> markdown helpers (previously imported from
# analyze_pptx.py). Kept here so this script is self-contained and portable.
# ---------------------------------------------------------------------------

def _chart_to_markdown(shape):
    if not getattr(shape, "has_chart", False):
        return []
    try:
        chart = shape.chart
    except Exception:
        return []
    lines = []
    ctype = str(chart.chart_type).split(".")[-1] if chart.chart_type is not None else "chart"
    chart_title = ""
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            chart_title = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass
    header = f"[chart: {ctype}]" + (f" {chart_title}" if chart_title else "")
    lines.append(header)
    categories = []
    try:
        plots = list(chart.plots)
        if plots:
            categories = [str(c) for c in plots[0].categories]
    except Exception:
        pass
    for series in chart.series:
        try:
            name = series.name or ""
        except Exception:
            name = ""
        try:
            values = [v for v in series.values]
        except Exception:
            values = []
        if categories and len(categories) == len(values):
            pairs = ", ".join(f"{cat}: {val}" for cat, val in zip(categories, values))
            lines.append(f"- {name}: {pairs}" if name else f"- {pairs}")
        elif values:
            lines.append(f"- {name}: {values}" if name else f"- {values}")
        elif name:
            lines.append(f"- {name}")
    return lines


def _picture_to_markdown(shape):
    name = (getattr(shape, "name", "") or "").strip()
    alt = ""
    try:
        cnv = shape._element.xpath(".//p:nvPicPr/p:cNvPr")
        if cnv:
            alt = (cnv[0].get("descr") or cnv[0].get("title") or "").strip()
    except Exception:
        pass
    fname = ""
    try:
        fname = (shape.image.filename or "").strip()
    except Exception:
        pass
    bits = [b for b in [alt, name, fname] if b and b.lower() not in {"picture", "image"}]
    return [f"[image] {' — '.join(dict.fromkeys(bits))}"] if bits else ["[image]"]


def _smartart_to_markdown(shape):
    try:
        texts = shape._element.xpath(".//a:t")
    except Exception:
        return []
    out = [f"- {(t.text or '').strip()}" for t in texts if (t.text or "").strip()]
    return (["[smartart]"] + out) if out else []


def _shape_first_line(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    for para in shape.text_frame.paragraphs:
        t = "".join(run.text for run in para.runs).strip()
        if t:
            return t
    return ""


def _infer_title(slide):
    if slide.shapes.title is not None:
        t = (slide.shapes.title.text or "").strip()
        if t:
            return t
    candidates = []
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        line = _shape_first_line(shape)
        if not line:
            continue
        top = shape.top if shape.top is not None else 10**9
        candidates.append((top, line))
    if not candidates:
        return ""
    candidates.sort(key=lambda x: x[0])
    title = candidates[0][1]
    return title if len(title) <= 120 else title[:117] + "…"


def _shape_to_markdown(shape, title_shape):
    if shape == title_shape:
        return []
    lines = []
    stype = getattr(shape, "shape_type", None)
    if stype == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            lines.extend(_shape_to_markdown(sub, title_shape))
        return lines
    if getattr(shape, "has_chart", False):
        lines.extend(_chart_to_markdown(shape))
        if lines:
            lines.append("")
        return lines
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
        return lines
    if stype == MSO_SHAPE_TYPE.PICTURE:
        lines.extend(_picture_to_markdown(shape))
        return lines
    if stype in (getattr(MSO_SHAPE_TYPE, "DIAGRAM", None), 15):
        sm = _smartart_to_markdown(shape)
        if sm:
            lines.extend(sm); lines.append("")
            return lines
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                prefix = "  " * (para.level or 0) + "- "
                lines.append(prefix + t)
        return lines
    try:
        leftover = [t.text.strip() for t in shape._element.xpath(".//a:t")
                    if t.text and t.text.strip()]
        if leftover:
            lines.append("- " + " ".join(leftover))
    except Exception:
        pass
    return lines


def _notes_to_markdown(slide):
    if not slide.has_notes_slide:
        return []
    nf = slide.notes_slide.notes_text_frame
    if nf is None:
        return []
    notes = []
    for para in nf.paragraphs:
        t = "".join(run.text for run in para.runs).strip()
        if t:
            notes.append(t)
    if not notes:
        return []
    return ["", "## Speaker notes", *[f"> {n}" for n in notes]]


def _layout_name(slide):
    try:
        return (slide.slide_layout.name or "").strip()
    except Exception:
        return ""


def _slide_to_markdown(slide):
    title_text = _infer_title(slide)
    title_shape = slide.shapes.title
    body_lines = []
    for shape in slide.shapes:
        body_lines.extend(_shape_to_markdown(shape, title_shape))
    body_lines.extend(_notes_to_markdown(slide))
    parts = []
    layout = _layout_name(slide)
    header_bits = []
    if title_text:
        header_bits.append(f"# {title_text}")
    if layout:
        header_bits.append(f"_layout: {layout}_")
    if header_bits:
        parts.append("\n".join(header_bits))
    body = "\n".join(body_lines).strip()
    if body:
        parts.append(body)
    return title_text, "\n\n".join(parts).strip()

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

load_dotenv()

# Foundry / Azure OpenAI endpoint hosting the vision deployment.
AOAI_ENDPOINT = (
    os.environ.get("AOAI_ENDPOINT")
    or os.environ.get("FOUNDRY_ENDPOINT")
    or ""
).rstrip("/")
AOAI_API_VER  = os.environ.get("AOAI_API_VERSION", "2024-10-21")
VISION_DEPLOY = os.environ.get("VISION_DEPLOYMENT", "gpt-5.4")

DEFAULT_DPI = int(os.environ.get("RENDER_DPI", "200"))

if not AOAI_ENDPOINT:
    print("ERROR: set AOAI_ENDPOINT (or FOUNDRY_ENDPOINT) in .env",
          file=sys.stderr)
    sys.exit(2)

credential = ChainedTokenCredential(
    AzureCliCredential(process_timeout=60),
    DefaultAzureCredential(exclude_interactive_browser_credential=False),
)
aoai_token_provider = get_bearer_token_provider(
    credential, "https://cognitiveservices.azure.com/.default"
)
# Per-request timeout (seconds). High-detail vision calls can be slow but
# should never hang forever — without this, a stuck TLS read blocks the
# whole pipeline. Override with VISION_TIMEOUT env var.
VISION_TIMEOUT = float(os.environ.get("VISION_TIMEOUT", "120"))

aoai = AzureOpenAI(
    azure_endpoint=AOAI_ENDPOINT,
    azure_ad_token_provider=aoai_token_provider,
    api_version=AOAI_API_VER,
    timeout=VISION_TIMEOUT,
    max_retries=0,  # we handle retries ourselves in extract_slide()
)

# ---------------------------------------------------------------------------
# JSON schema sent to the vision model (mirrors FIELDS_DEF in analyze_pptx.py)
# ---------------------------------------------------------------------------

SLIDE_JSON_SCHEMA = {
    "name": "SlideFields",
    "strict": False,
    "schema": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "slide_title":       {"type": "string"},
            "slide_summary":     {"type": "string"},
            "slide_description": {"type": "string"},
            "highlighted_terms": {
                "type": "array",
                "description": "Words/phrases the slide visually emphasises (pink/red/yellow fill, colored underline, bold-colored text). Each entry is the verbatim emphasised text with optional color/location.",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "text":     {"type": "string"},
                        "color":    {"type": "string"},
                        "location": {"type": "string", "description": "e.g. title, subtitle, body, callout, footer"},
                    },
                    "required": ["text", "color", "location"],
                },
            },
            "kpis": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "label": {"type": "string"},
                        "value": {"type": "string"},
                        "delta": {"type": "string"},
                    },
                    "required": ["label", "value", "delta"],
                },
            },
            "charts": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "title":      {"type": "string"},
                        "chart_type": {"type": "string"},
                        "x_axis":     {"type": "string"},
                        "y_axis":     {"type": "string"},
                        "unit":       {"type": "string"},
                        "categories": {"type": "array", "items": {"type": "string"}},
                        "series": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "additionalProperties": False,
                                "properties": {
                                    "name":   {"type": "string"},
                                    "color":  {"type": "string"},
                                    "values": {"type": "array", "items": {"type": "string"}},
                                },
                                "required": ["name", "color", "values"],
                            },
                        },
                        "highlights": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "additionalProperties": False,
                                "properties": {
                                    "category": {"type": "string"},
                                    "series":   {"type": "string"},
                                    "value":    {"type": "string"},
                                    "note":     {"type": "string"},
                                },
                                "required": ["category", "series", "value", "note"],
                            },
                        },
                        "annotations": {
                            "type": "array",
                            "description": "Visual overlays such as dashed/solid boxes, arrows, numbered callouts.",
                            "items": {
                                "type": "object",
                                "additionalProperties": False,
                                "properties": {
                                    "label":    {"type": "string", "description": "Number/letter inside the callout, if any."},
                                    "shape":    {"type": "string", "description": "dashed-box, solid-box, arrow, circle, callout, etc."},
                                    "color":    {"type": "string"},
                                    "encloses": {"type": "string", "description": "What the overlay points at or contains, in plain English."},
                                    "meaning":  {"type": "string", "description": "Why it is highlighted, if inferable from the slide."},
                                },
                                "required": ["label", "shape", "color", "encloses", "meaning"],
                            },
                        },
                        "source_table_index": {"type": "string"},
                        "takeaway":   {"type": "string"},
                    },
                    "required": ["title", "chart_type", "categories", "series",
                                 "highlights", "takeaway"],
                },
            },
            "tables": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "title":   {"type": "string"},
                        "headers": {"type": "array", "items": {"type": "string"}},
                        "rows": {
                            "type": "array",
                            "items": {"type": "array", "items": {"type": "string"}},
                        },
                        "row_label_column": {"type": "string"},
                        "cell_highlights": {
                            "type": "array",
                            "description": "Cells that are visually emphasised (colored fill, bold border, etc).",
                            "items": {
                                "type": "object",
                                "additionalProperties": False,
                                "properties": {
                                    "row_label":    {"type": "string"},
                                    "column_label": {"type": "string"},
                                    "value":        {"type": "string"},
                                    "color":        {"type": "string"},
                                    "meaning":      {"type": "string"},
                                },
                                "required": ["row_label", "column_label", "value", "color", "meaning"],
                            },
                        },
                        "notes": {"type": "string"},
                    },
                    "required": ["headers", "rows"],
                },
            },
            "callouts": {
                "type": "array",
                "description": "Standalone text callout boxes / pull-quotes on the slide.",
                "items": {"type": "string"},
            },
        },
        "required": ["slide_title", "slide_summary", "slide_description",
                     "kpis", "charts", "tables"],
    },
}

EXTRACT_SYSTEM = (
    "You are a meticulous slide-reading assistant. You receive ONE slide image "
    "from a corporate deck plus ground-truth markdown extracted from the "
    "underlying PowerPoint XML. Use the image for layout, colors, callouts, "
    "annotations, and visual emphasis. Use the markdown for exact numbers when "
    "it contradicts what you might OCR. Return STRICT JSON matching the "
    "provided schema. Do not invent numbers that are not visible or in the "
    "ground truth. "
    "Pay special attention to visually highlighted words/phrases (pink, red, "
    "yellow, or other colored fill / underline / bold-color text). Capture "
    "every such emphasis in `highlighted_terms`, and ALSO mention them "
    "verbatim in `slide_description`. If the slide title contains a "
    "highlighted word, append it at the end of `slide_title` as "
    "`[highlighted: <word>]` so downstream search picks it up. "
    "If the slide has NO visual highlights, return an empty list `[]` for "
    "`highlighted_terms` and do NOT append any `[highlighted: …]` suffix to "
    "the title — never invent emphasis that isn't visibly there."
)

EXTRACT_USER_TEMPLATE = """Slide {slide_number} of {total_slides}.

== GROUND-TRUTH MARKDOWN (from PPTX XML) ==
{markdown}
== END GROUND-TRUTH ==

Fill the schema. Specifically:
- `slide_title`: full visible title. If any word/phrase is visually highlighted
  (pink/red/yellow fill, colored underline, bold-color), append it at the end
  like `… [highlighted: Nation]` (multiple → comma-separated).
- `slide_description`: 1–3 sentences. Explicitly mention every highlighted
  word/phrase verbatim and what it emphasizes.
- `highlighted_terms`: list every visually emphasised word/phrase on the slide
  (titles, body, callouts) with its color and location. Do NOT skip any.
  If there are NO visual highlights, return `[]` (empty list) and do not
  append any `[highlighted: …]` suffix to `slide_title`.
- `charts`: capture every series with values aligned to `categories` (same length, same order).
- `tables`: capture the full grid; if a chart visualises a table on this slide,
  set its `source_table_index` to the table's 0-based index.
- `highlights` (per chart) and `cell_highlights` (per table): record peaks,
  dips, color-shaded cells, arrows, +/- deltas the slide is drawing attention to.
- `annotations` (per chart): record dashed/solid overlay boxes, numbered "1"/"2"
  callouts, arrows — anything visual that is NOT part of the chart data.
- `callouts`: pull-quotes / standalone text boxes that are not chart titles.

Return JSON only.
"""

# ---------------------------------------------------------------------------
# Render helpers
# ---------------------------------------------------------------------------

def _ensure_pdf(pptx: Path, pdf: Path | None, expected_pages: int | None = None) -> Path:
    if pdf and pdf.exists():
        return pdf
    sibling = pptx.with_suffix(".pdf")
    if sibling.exists():
        # Detect stale PDF (page count != slide count). If LibreOffice is
        # available we re-render; otherwise warn and continue with the stale one.
        if expected_pages is not None:
            try:
                with fitz.open(sibling) as _doc:
                    actual = _doc.page_count
            except Exception:
                actual = None
            if actual is not None and actual != expected_pages:
                soffice = shutil.which("soffice") or shutil.which("libreoffice")
                if soffice:
                    print(f"  stale PDF ({actual} pages vs {expected_pages} slides) — "
                          f"re-rendering via {soffice}…")
                    subprocess.run(
                        [soffice, "--headless", "--convert-to", "pdf",
                         "--outdir", str(pptx.parent), str(pptx)],
                        check=True,
                    )
                else:
                    print(f"  WARNING: PDF has {actual} pages but pptx has "
                          f"{expected_pages} slides. Re-export the PDF manually "
                          "to match, or install LibreOffice for auto-rendering.")
        return sibling
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError(
            f"No PDF found at {sibling} and LibreOffice (`soffice`) is not on "
            "PATH. Either export the deck to PDF manually next to the .pptx, "
            "or install LibreOffice and re-run."
        )
    print(f"  converting {pptx.name} → PDF via {soffice}…")
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf",
         "--outdir", str(pptx.parent), str(pptx)],
        check=True,
    )
    if not sibling.exists():
        raise RuntimeError(f"LibreOffice did not produce {sibling}")
    return sibling


def _render_pdf_to_pngs(pdf: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    paths: list[Path] = []
    with fitz.open(pdf) as doc:
        for i, page in enumerate(doc, 1):
            out = out_dir / f"slide_{i:02d}.png"
            page.get_pixmap(matrix=matrix, alpha=False).save(str(out))
            paths.append(out)
    return paths


def _b64_png(path: Path) -> str:
    return base64.b64encode(path.read_bytes()).decode("ascii")


# ---------------------------------------------------------------------------
# Vision call
# ---------------------------------------------------------------------------

def extract_slide(slide_number: int, total: int, img: Path,
                  ground_truth_md: str, *, max_retries: int = 2) -> dict:
    user_text = EXTRACT_USER_TEMPLATE.format(
        slide_number=slide_number,
        total_slides=total,
        markdown=ground_truth_md or "(no underlying text/chart data extracted)",
    )
    last_err: Exception | None = None
    for attempt in range(max_retries + 1):
        try:
            resp = aoai.chat.completions.create(
                model=VISION_DEPLOY,
                response_format={"type": "json_schema",
                                 "json_schema": SLIDE_JSON_SCHEMA},
                temperature=0,
                messages=[
                    {"role": "system", "content": EXTRACT_SYSTEM},
                    {"role": "user", "content": [
                        {"type": "text", "text": user_text},
                        {"type": "image_url", "image_url": {
                            "url": f"data:image/png;base64,{_b64_png(img)}",
                            "detail": "high",
                        }},
                    ]},
                ],
            )
            return json.loads(resp.choices[0].message.content)
        except Exception as e:  # noqa: BLE001
            last_err = e
            if attempt < max_retries:
                wait = 2 * (attempt + 1)
                print(f"    retry {attempt + 1} after error ({e}); sleeping {wait}s")
                time.sleep(wait)
    raise RuntimeError(f"Vision call failed after {max_retries + 1} attempts: {last_err}")


# ---------------------------------------------------------------------------
# Flatten helpers (re-implemented locally to include vision-only fields)
# ---------------------------------------------------------------------------

def _flatten_kpis(kpis):
    parts = []
    for k in kpis or []:
        label, value, delta = k.get("label", ""), k.get("value", ""), k.get("delta", "")
        if not (label or value):
            continue
        s = f"{label}: {value}".strip(": ").strip()
        if delta:
            s += f" ({delta})"
        parts.append(s)
    return "; ".join(parts)


def _flatten_charts(charts):
    parts = []
    for ch in charts or []:
        head = f"[{ch.get('chart_type','')}] {ch.get('title','')}".strip("[] ").strip()
        bits = [head] if head else []
        cats = ch.get("categories") or []
        for s in ch.get("series", []) or []:
            vals = s.get("values") or []
            pairs = [f"{c}: {v}" for c, v in zip(cats, vals) if (c or v)]
            if pairs:
                name = s.get("name") or ""
                bits.append(f"{name} -> {', '.join(pairs)}" if name else ", ".join(pairs))
            elif s.get("name"):
                bits.append(s["name"])
        for h in ch.get("highlights", []) or []:
            label = "/".join(x for x in [h.get("series"), h.get("category")] if x)
            seg = f"★ {label}: {h.get('value','')}".strip(": ").strip()
            if h.get("note"):
                seg = f"{seg} ({h['note']})" if seg else h["note"]
            if seg:
                bits.append(seg)
        for a in ch.get("annotations", []) or []:
            descr = " ".join(x for x in [a.get("shape"), a.get("color"), a.get("label")] if x).strip()
            if a.get("encloses"):
                descr = f"{descr} → {a['encloses']}".strip(" →")
            if a.get("meaning"):
                descr = f"{descr} ({a['meaning']})" if descr else a["meaning"]
            if descr:
                bits.append(f"⟦{descr}⟧")
        if ch.get("takeaway"):
            bits.append(ch["takeaway"])
        if bits:
            parts.append(" — ".join(bits))
    return " | ".join(parts)


def _flatten_tables(tables):
    parts = []
    for t in tables or []:
        head = t.get("title") or ""
        headers = t.get("headers") or []
        rows = t.get("rows") or []
        bits = [head] if head else []
        if headers:
            bits.append("cols: " + " | ".join(headers))
        for r in rows:
            bits.append(" | ".join(r))
        for h in t.get("cell_highlights", []) or []:
            seg = f"★ {h.get('row_label','')}/{h.get('column_label','')} = {h.get('value','')}"
            if h.get("color"):   seg += f" [{h['color']}]"
            if h.get("meaning"): seg += f" — {h['meaning']}"
            bits.append(seg)
        if t.get("notes"):
            bits.append(t["notes"])
        if bits:
            parts.append(" — ".join(bits))
    return " || ".join(parts)


# ---------------------------------------------------------------------------
# Slide range parsing
# ---------------------------------------------------------------------------

def _parse_slide_range(spec: str | None, total: int) -> set[int] | None:
    if not spec:
        return None
    keep: set[int] = set()
    for chunk in spec.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            a, b = chunk.split("-", 1)
            keep.update(range(int(a), int(b) + 1))
        else:
            keep.add(int(chunk))
    return {n for n in keep if 1 <= n <= total}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(description="Vision-based per-slide extraction for a .pptx")
    ap.add_argument("pptx", type=Path, help="Path to .pptx file")
    ap.add_argument("--out-dir", type=Path, default=None,
                    help="Output directory (default: same as input)")
    ap.add_argument("--pdf", type=Path, default=None,
                    help="Pre-rendered PDF (default: <deck>.pdf next to the pptx)")
    ap.add_argument("--img-dir", type=Path, default=None,
                    help="Directory to write rendered slide PNGs (default: ./slide_images_vision)")
    ap.add_argument("--dpi", type=int, default=DEFAULT_DPI,
                    help=f"Render DPI (default: {DEFAULT_DPI})")
    ap.add_argument("--only-slides", default=None,
                    help='Slide range to process, e.g. "9" or "1-3,9,12-14"')
    ap.add_argument("--skip-render", action="store_true",
                    help="Reuse PNGs already in --img-dir")
    args = ap.parse_args()

    pptx = args.pptx.resolve()
    if not pptx.exists():
        print(f"ERROR: file not found: {pptx}", file=sys.stderr)
        sys.exit(2)
    out_dir = (args.out_dir or pptx.parent).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    img_dir = (args.img_dir or (pptx.parent / "slide_images_vision")).resolve()

    print(f"Deck:     {pptx}  ({pptx.stat().st_size/1024:.1f} KB)")
    print(f"Endpoint: {AOAI_ENDPOINT}")
    print(f"Model:    {VISION_DEPLOY}  (api {AOAI_API_VER})")
    print(f"Out dir:  {out_dir}")
    print(f"Img dir:  {img_dir}\n")

    # 1. Render slides → PNG.
    if args.skip_render:
        img_paths = sorted(img_dir.glob("slide_*.png"))
        if not img_paths:
            print(f"ERROR: --skip-render set but no PNGs in {img_dir}", file=sys.stderr)
            sys.exit(2)
        print(f"[1/3] reusing {len(img_paths)} PNGs from {img_dir}")
    else:
        print("[1/3] rendering slides…")
        # Peek at slide count first so we can detect a stale sibling PDF.
        try:
            expected = len(Presentation(str(pptx)).slides)
        except Exception:
            expected = None
        pdf = _ensure_pdf(pptx, args.pdf, expected_pages=expected)
        img_paths = _render_pdf_to_pngs(pdf, img_dir, dpi=args.dpi)
        print(f"      wrote {len(img_paths)} PNGs to {img_dir}")

    # 2. Extract underlying markdown per slide via python-pptx (ground truth).
    print("\n[2/3] extracting ground-truth markdown via python-pptx…")
    prs = Presentation(str(pptx))
    slides_md: list[tuple[str, str]] = []  # (title, markdown)
    for slide in prs.slides:
        slides_md.append(_slide_to_markdown(slide))
    if len(slides_md) != len(img_paths):
        print(f"      WARNING: {len(slides_md)} pptx slides vs {len(img_paths)} PNGs; "
              "using min count")
    n = min(len(slides_md), len(img_paths))

    keep = _parse_slide_range(args.only_slides, n)
    if keep is not None:
        print(f"      restricted to slides: {sorted(keep)}")

    # 3. Vision call per slide.
    print(f"\n[3/3] calling {VISION_DEPLOY} per slide…")
    per_slide: list[dict] = []
    for i in range(1, n + 1):
        if keep is not None and i not in keep:
            continue
        title, md = slides_md[i - 1]
        img = img_paths[i - 1]
        print(f"  slide {i:>2}: {img.name} ({img.stat().st_size/1024:.0f} KB) "
              f"+ {len(md)} md chars …", end=" ", flush=True)
        try:
            data = extract_slide(i, n, img, md)
        except Exception as e:  # noqa: BLE001
            print(f"FAILED — {e}")
            continue
        # Normalise + envelope.
        data.setdefault("slide_title", title or "")
        data.setdefault("slide_summary", "")
        data.setdefault("slide_description", "")
        data.setdefault("highlighted_terms", [])
        data.setdefault("kpis", [])
        data.setdefault("charts", [])
        data.setdefault("tables", [])
        data.setdefault("callouts", [])
        # Coerce source_table_index to int|None for ergonomics.
        for ch in data.get("charts", []):
            raw = ch.get("source_table_index", "")
            try:
                ch["source_table_index"] = int(raw) if raw not in ("", None) else None
            except (TypeError, ValueError):
                ch["source_table_index"] = None
        per_slide.append({
            "slide_number":      i,
            "page_number":       i,
            "slide_title":       data["slide_title"],
            "slide_summary":     data["slide_summary"],
            "slide_description": data["slide_description"],
            "highlighted_terms": data["highlighted_terms"],
            "kpis":              data["kpis"],
            "charts":            data["charts"],
            "tables":            data["tables"],
            "callouts":          data["callouts"],
            "markdown":          md,
            "image":             str(img.relative_to(img.parent.parent))
                                 if img.is_relative_to(img.parent.parent) else img.name,
        })
        print(f"ok — kpis={len(data['kpis'])}, "
              f"charts={len(data['charts'])}, "
              f"tables={len(data['tables'])}, "
              f"callouts={len(data['callouts'])}")

    # 4. Persist.
    deck_stem = pptx.stem
    rich_path  = out_dir / f"{deck_stem}.vision.slides.json"
    jsonl_path = out_dir / f"{deck_stem}.vision.slides.jsonl"

    rich_path.write_text(
        json.dumps({"deck": pptx.name, "model": VISION_DEPLOY,
                    "slides": per_slide}, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    search_docs = []
    for s in per_slide:
        callouts_text = " || ".join(s.get("callouts", []))
        highlights_text = " || ".join(
            f"{h.get('text','')} [{h.get('color','')} @ {h.get('location','')}]".strip()
            for h in (s.get("highlighted_terms") or [])
            if h.get("text")
        )
        search_docs.append({
            "id":                f"{deck_stem}_slide_{s['slide_number']:02d}",
            "deck":              pptx.name,
            "slide_number":      s["slide_number"],
            "page_number":       s["page_number"],
            "slide_title":       s["slide_title"],
            "slide_summary":     s["slide_summary"],
            "slide_description": s["slide_description"],
            "highlighted_terms": highlights_text,
            "kpis_text":         _flatten_kpis(s["kpis"]),
            "charts_text":       _flatten_charts(s["charts"]),
            "tables_text":       _flatten_tables(s["tables"]),
            "callouts_text":     callouts_text,
            "content":           s["markdown"],
        })
    with jsonl_path.open("w", encoding="utf-8") as f:
        for d in search_docs:
            f.write(json.dumps(d, ensure_ascii=False) + "\n")

    print(f"\nWrote {rich_path}")
    print(f"Wrote {jsonl_path}  ({len(search_docs)} docs)")


if __name__ == "__main__":
    main()
